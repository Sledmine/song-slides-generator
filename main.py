from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def parse_blocks(filename):
    with open(filename, 'r', encoding='utf-8') as f:
        lines = [line.rstrip('\n') for line in f]

    blocks = []
    current_block = []

    for line in lines:
        if line.strip() == "":
            if current_block:
                blocks.append(current_block)
                current_block = []
        else:
            current_block.append(line)

    if current_block:
        blocks.append(current_block)

    return blocks

def create_ppt(blocks, output_file="output.pptx"):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # blank slide

    for block in blocks:
        slide = prs.slides.add_slide(blank_slide_layout)
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(5.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.clear()

        # Each pair of lines: chord line followed by lyrics line
        i = 0
        while i < len(block):
            chords = block[i]
            lyrics = block[i + 1] if i + 1 < len(block) else ""

            # Chord line
            chord_p = tf.add_paragraph()
            chord_p.space_after = Pt(0)
            chord_run = chord_p.add_run()
            chord_run.text = chords
            chord_run.font.bold = True
            chord_run.font.size = Pt(20)
            #chord_run.font.color.rgb = RGBColor(255, 0, 0) # Red
            chord_run.font.color.rgb = RGBColor(0x42, 0x24, 0xE9)

            # Lyrics line
            lyric_p = tf.add_paragraph()
            lyric_p.space_after = Pt(10)
            lyric_run = lyric_p.add_run()
            lyric_run.text = lyrics
            lyric_run.font.size = Pt(20)

            i += 2  # Move to next pair

    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

blocks = parse_blocks("song.txt")
create_ppt(blocks)
